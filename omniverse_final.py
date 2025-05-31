import customtkinter as ctk
from PIL import Image
import pyttsx3
import sounddevice as sd
import queue
import vosk
import sys
import os
import json
import threading
import google.generativeai as genai
from gtts import gTTS
from pydub import AudioSegment
from diffusers import StableDiffusionPipeline
import torch
import datetime
from fpdf import FPDF
import docx
from pptx import Presentation
import pandas as pd

# === Gemini API ===
genai.configure(api_key="your_api_key")
model = genai.GenerativeModel("gemini-1.5-flash")

# === Vosk Model ===
model_path = "vosk-model-small-en-us-0.15"
if not os.path.exists(model_path):
    print("Download Vosk model and unzip as 'vosk-model-small-en-us-0.15'")
    sys.exit()
vosk_model = vosk.Model(model_path)
q = queue.Queue()

# === TTS Engine ===
engine = pyttsx3.init()
is_speaking = False


def speak_text(text):
    global is_speaking
    is_speaking = True
    engine.say(text)
    try:
        engine.runAndWait()
    except RuntimeError:
        pass
    is_speaking = False


def stop_speaking():
    global is_speaking
    if is_speaking:
        engine.stop()
        chat_display.insert(ctk.END, "\n🛑 AI speaking stopped.\n")
        chat_display.see(ctk.END)
        is_speaking = False


# === Image Generator ===
def generate_image(prompt, output_path="output.png"):
    pipe = StableDiffusionPipeline.from_pretrained("runwayml/stable-diffusion-v1-5", torch_dtype=torch.float32)
    pipe = pipe.to("cpu")
    image = pipe(prompt).images[0]
    image.save(output_path)
    chat_display.insert(ctk.END, f"🖼️ Image saved as {output_path}\n")


# === Singing ===
def sing_lyrics(lyrics, lang='en'):
    tts = gTTS(text=lyrics, lang=lang)
    tts.save("song.mp3")
    song = AudioSegment.from_mp3("song.mp3")
    song.export("final_song.wav", format="wav")
    os.system("start final_song.wav")
    chat_display.insert(ctk.END, "🎵 Singing your lyrics...\n")


# === Video Simulation ===
def simulate_video_creation(prompt):
    generate_image(prompt, "scene.png")
    sing_lyrics("Narrating: " + prompt)
    chat_display.insert(ctk.END, "🎬 Video simulation complete.\n")


# === Resume Builder ===
def build_resume():
    user_info = entry.get()
    if not user_info.strip():
        chat_display.insert(ctk.END, "⚠️ Please enter role or key skills to generate a resume.\n")
        return
    try:
        chat_display.insert(ctk.END, "🧾 Creating ATS-friendly resume...\n")
        chat_display.see(ctk.END)
        prompt = (
            f"Create a professional ATS-friendly resume for a candidate skilled in {user_info}. "
            "Include Name, Summary, Skills, Work Experience, Education, Certifications, and Projects. "
            "Keep formatting plain and machine-readable."
        )
        response = model.generate_content(prompt)
        resume_text = response.text.strip()
        filename = f"ATS_Resume_{datetime.datetime.now().strftime('%Y%m%d%H%M%S')}.txt"
        with open(filename, "w", encoding="utf-8") as file:
            file.write(resume_text)
        chat_display.insert(ctk.END, f"✅ Resume generated: {filename}\n")
        threading.Thread(target=speak_text, args=("Your resume is ready.",), daemon=True).start()
    except Exception as e:
        chat_display.insert(ctk.END, f"❌ Error generating resume: {e}\n")


# === Health Prediction ===
def predict_health_risk():
    user_input = entry.get()
    if not user_input.strip():
        chat_display.insert(ctk.END, "⚠️ Please enter symptoms or health indicators.\n")
        return
    try:
        prompt = (
            f"Based on these health symptoms or indicators: {user_input}, "
            "provide an estimated health risk (e.g., diabetes, heart disease) and a confidence level above 90%. "
            "Keep it short and expert-like."
        )
        chat_display.insert(ctk.END, "🩺 Predicting health status...\n")
        chat_display.see(ctk.END)
        response = model.generate_content(prompt)
        result = response.text.strip()
        chat_display.insert(ctk.END, f"🩺 Health Prediction Result: {result}\n")
        chat_display.see(ctk.END)
        threading.Thread(target=speak_text, args=(result,), daemon=True).start()
    except Exception as e:
        chat_display.insert(ctk.END, f"❌ Health Prediction Error: {e}\n")


# === PDF, DOCX, PPTX, EXCEL Creator ===
def create_all_documents():
    user_input = entry.get()
    if not user_input.strip():
        chat_display.insert(ctk.END, "⚠️ Please enter content to generate documents.\n")
        return
    try:
        response = model.generate_content(user_input)
        content = response.text.strip()

        # PDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        for line in content.split('\n'):
            pdf.multi_cell(0, 10, line)
        pdf_file = "generated_document.pdf"
        pdf.output(pdf_file)

        # Word
        doc = docx.Document()
        doc.add_paragraph(content)
        docx_file = "generated_document.docx"
        doc.save(docx_file)

        # PPT
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = "AI Generated Presentation"
        slide.placeholders[1].text = content[:500]  # limit to fit
        pptx_file = "generated_document.pptx"
        ppt.save(pptx_file)

        # Excel
        df = pd.DataFrame({'Content': content.split('\n')})
        excel_file = "generated_document.xlsx"
        df.to_excel(excel_file, index=False)

        chat_display.insert(ctk.END, f"✅ Documents created: {pdf_file}, {docx_file}, {pptx_file}, {excel_file}\n")
    except Exception as e:
        chat_display.insert(ctk.END, f"❌ Error generating documents: {e}\n")


# === GUI Setup ===
ctk.set_appearance_mode("dark")
ctk.set_default_color_theme("blue")
app = ctk.CTk()
app.geometry("720x980")
app.title("Omniverse AI Assistant")

# === Logo ===
logo_image = Image.open("D:/omniverse_ai/Omniverse AI Logo Design.png").resize((100, 100))
logo = ctk.CTkImage(dark_image=logo_image, light_image=logo_image, size=(100, 100))
ctk.CTkLabel(app, image=logo, text="").pack(pady=10)

# === Chat Display ===
chat_display = ctk.CTkTextbox(app, width=680, height=300, font=("Arial", 14))
chat_display.pack(pady=10)
chat_display.insert(ctk.END, "Welcome to Omniverse AI!\n\n")

# === Entry ===
entry = ctk.CTkEntry(app, width=680, font=("Arial", 14))
entry.pack(pady=10)

# === Core Functions ===
def generate_response(user_input):
    try:
        chat_display.insert(ctk.END, "Omniverse AI: Thinking...\n")
        chat_display.see(ctk.END)
        response = model.generate_content(user_input)
        reply = response.text.strip()
        chat_display.insert(ctk.END, f"Omniverse AI: {reply}\n")
        chat_display.see(ctk.END)
        threading.Thread(target=speak_text, args=(reply,), daemon=True).start()
    except Exception as e:
        chat_display.insert(ctk.END, f"❌ Error: {e}\n")


def send_input():
    user_text = entry.get()
    if user_text.strip():
        entry.delete(0, ctk.END)
        chat_display.insert(ctk.END, f"You: {user_text}\n")
        chat_display.see(ctk.END)
        threading.Thread(target=generate_response, args=(user_text,), daemon=True).start()


def callback(indata, frames, time, status):
    q.put(bytes(indata))


def recognize_speech():
    rec = vosk.KaldiRecognizer(vosk_model, 16000)
    with sd.RawInputStream(samplerate=16000, blocksize=8000, dtype='int16', channels=1, callback=callback):
        print("Listening...")
        while True:
            data = q.get()
            if rec.AcceptWaveform(data):
                result = json.loads(rec.Result())
                text = result.get("text", "")
                if text:
                    chat_display.insert(ctk.END, f"You (voice): {text}\n")
                    chat_display.see(ctk.END)
                    threading.Thread(target=generate_response, args=(text,), daemon=True).start()
                    break

# === Buttons ===
scrollable_frame = ctk.CTkScrollableFrame(app, width=700, height=500)
scrollable_frame.pack(pady=10)

ctk.CTkButton(scrollable_frame, text="💬 Send", width=300, command=lambda: threading.Thread(target=send_input, daemon=True).start()).pack(pady=5)
ctk.CTkButton(scrollable_frame, text="🎙️ Speak", width=300, command=lambda: threading.Thread(target=recognize_speech, daemon=True).start()).pack(pady=5)
ctk.CTkButton(scrollable_frame, text="🛑 Stop Speaking", width=300, command=stop_speaking).pack(pady=5)
ctk.CTkButton(scrollable_frame, text="🖼️ Generate Image", width=300, command=lambda: threading.Thread(target=lambda: generate_image(entry.get()), daemon=True).start()).pack(pady=5)
ctk.CTkButton(scrollable_frame, text="🎵 Sing Lyrics", width=300, command=lambda: threading.Thread(target=lambda: sing_lyrics(entry.get()), daemon=True).start()).pack(pady=5)
ctk.CTkButton(scrollable_frame, text="🎬 Create Video", width=300, command=lambda: threading.Thread(target=lambda: simulate_video_creation(entry.get()), daemon=True).start()).pack(pady=5)
ctk.CTkButton(scrollable_frame, text="🩺 Predict Health", width=300, command=lambda: threading.Thread(target=predict_health_risk, daemon=True).start()).pack(pady=5)
ctk.CTkButton(scrollable_frame, text="📑 Build Resume", width=300, command=lambda: threading.Thread(target=build_resume, daemon=True).start()).pack(pady=5)
ctk.CTkButton(scrollable_frame, text="📁 Create All Documents", width=300, command=lambda: threading.Thread(target=create_all_documents, daemon=True).start()).pack(pady=5)

# === Start App ===
app.mainloop()
